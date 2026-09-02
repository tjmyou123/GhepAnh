"""Trinh chieu anh tu dong voi hieu ung zoom (Ken Burns) — khong can Markdown.

Chi can chon thu muc anh:
- HTML: moi anh phu kin man hinh, zoom in/out + luot cham, tu chuyen slide,
  crossfade, dieu khien Space/←/→/F, mot file duy nhat chay offline.
- PPTX: moi anh mot slide (cat vua khit khung), hieu ung zoom that (Grow/Shrink)
  + chuyen canh fade + tu dong chuyen sau N giay ngay trong PowerPoint.
"""

from __future__ import annotations

import random as _random
from pathlib import Path
from typing import Callable, Optional

from PIL import Image, ImageOps

from .core import CollageError, make_collage_image, unique_path
from .renderer import find_images, probe_images
from .slides_assets import data_uri, image_bytes_for_html, pil_to_stream
from .slides_themes import FONT_STACK, get_slide_theme
from .themes import THEMES

MAX_PHOTOS = 300


def _collect_photos(folder: str | Path, order: str = "name"
                    ) -> tuple[Path, list[Path], list[str]]:
    folder = Path(folder)
    if not folder.is_dir():
        raise CollageError(f"Khong tim thay thu muc: {folder}")
    files = find_images(folder)
    if not files:
        raise CollageError(
            "Thu muc khong co anh nao (ho tro: jpg, png, webp, bmp, tiff, gif).")
    warnings: list[str] = []
    if len(files) > MAX_PHOTOS:
        warnings.append(
            f"Thu muc co {len(files)} anh, chi dung {MAX_PHOTOS} anh dau tien.")
        files = files[:MAX_PHOTOS]
    files, _aspects, errors = probe_images(files)
    warnings.extend(f"Bo qua anh loi: {e}" for e in errors)
    if not files:
        raise CollageError("Tat ca anh trong thu muc deu bi loi.")
    if order == "random":
        _random.shuffle(files)
    if len(files) > 100:
        warnings.append(
            f"{len(files)} anh -> file xuat ra se kha nang, nen bot anh "
            "hoac chia nhieu phan.")
    return folder, files, warnings


def _intro_collage(folder: Path, theme: str, size: str,
                   status: Callable[[str], None]) -> Optional[Image.Image]:
    """Ghep collage mo dau tu toan bo anh (None neu loi)."""
    try:
        status("Đang ghép collage mở đầu...")
        img, _w = make_collage_image(
            folder, preset="ppt" if size == "16:9" else "ppt-43",
            layout_style="justified",
            theme=theme if theme in THEMES else "classic",
            supersample=1)
        return img
    except CollageError:
        return None


# ================================================================== HTML ====
_KB_CLASSES = ("k0", "k1", "k2", "k3", "k4", "k5")


def _slideshow_css(theme: dict, dur: float, zoom: int, size: str) -> str:
    z = 1 + max(4, min(40, zoom)) / 100          # 1.12 mac dinh
    p = 1 + max(4, min(40, zoom)) / 200 + 0.04   # muc phong khi luot ngang
    anim_dur = dur + 1.4
    return f"""
*{{box-sizing:border-box}}
html,body{{margin:0;height:100%;background:#000;overflow:hidden;
  font-family:{FONT_STACK}}}
.ph{{position:fixed;inset:0;opacity:0;transition:opacity 1.15s ease;
  overflow:hidden;z-index:1}}
.ph.active{{opacity:1;z-index:2}}
.ph img{{width:100%;height:100%;object-fit:cover;will-change:transform}}
.ph.active.k0 img{{animation:kb-in {anim_dur}s ease-in-out both}}
.ph.active.k1 img{{animation:kb-out {anim_dur}s ease-in-out both}}
.ph.active.k2 img{{animation:kb-tl {anim_dur}s ease-in-out both}}
.ph.active.k3 img{{animation:kb-br {anim_dur}s ease-in-out both}}
.ph.active.k4 img{{animation:kb-pan {anim_dur}s ease-in-out both}}
.ph.active.k5 img{{animation:kb-pan2 {anim_dur}s ease-in-out both}}
body.paused .ph img{{animation-play-state:paused}}
@keyframes kb-in{{from{{transform:scale(1)}}to{{transform:scale({z})}}}}
@keyframes kb-out{{from{{transform:scale({z})}}to{{transform:scale(1)}}}}
@keyframes kb-tl{{from{{transform:scale(1.03)}}to{{transform:scale({z}) translate(2.4%,1.7%)}}}}
@keyframes kb-br{{from{{transform:scale(1.03)}}to{{transform:scale({z}) translate(-2.4%,-1.7%)}}}}
@keyframes kb-pan{{from{{transform:scale({p}) translate(2.2%,0)}}to{{transform:scale({p}) translate(-2.2%,0)}}}}
@keyframes kb-pan2{{from{{transform:scale({p}) translate(-2.2%,0)}}to{{transform:scale({p}) translate(2.2%,0)}}}}
.ph.title{{display:grid;place-items:center;text-align:center;
  background:linear-gradient(160deg,{theme['bg'][0]},{theme['bg'][1]})}}
.ph.title .tt{{max-width:80%}}
.ph.title h1{{color:{theme['h1']};font-size:min(7vw,84px);font-weight:800;
  margin:0 0 18px;line-height:1.15}}
.ph.title p{{color:{theme['muted']};font-size:min(2.6vw,26px);margin:6px 0}}
.ph.title .hint{{color:{theme['accent']};margin-top:34px;font-size:17px}}
.cap{{position:fixed;left:22px;bottom:18px;z-index:10;color:#EEE;font-size:15px;
  background:rgba(0,0,0,.45);padding:6px 14px;border-radius:9px;
  backdrop-filter:blur(6px);max-width:60vw;overflow:hidden;
  text-overflow:ellipsis;white-space:nowrap}}
#bar{{position:fixed;left:0;top:0;height:4px;width:0;z-index:20;
  background:linear-gradient(90deg,{theme['accent']},{theme['h2']})}}
#bar.run{{animation:pbar {dur}s linear both}}
body.paused #bar.run{{animation-play-state:paused}}
@keyframes pbar{{from{{width:0}}to{{width:100%}}}}
#ui{{position:fixed;right:16px;bottom:14px;z-index:30;display:flex;gap:7px;
  align-items:center;opacity:.95;transition:opacity .4s}}
body.idle #ui{{opacity:0;pointer-events:none}}
body.idle{{cursor:none}}
#ui button{{border:none;border-radius:10px;background:rgba(25,25,32,.72);
  color:#EEE;width:42px;height:37px;font-size:16px;cursor:pointer;
  backdrop-filter:blur(6px)}}
#ui button:hover{{background:rgba(70,70,90,.9)}}
#cnt{{color:#CCC;font-size:13px;padding:0 7px;user-select:none}}
"""


_SHOW_JS = r"""
(function(){
var phs=[].slice.call(document.querySelectorAll('.ph')),N=phs.length,cur=0,
 playing=true,t0=0,left=0,timer=null,DUR=window.__DUR*1000,LOOP=window.__LOOP,
 bar=document.getElementById('bar'),cnt=document.getElementById('cnt'),
 btnP=document.getElementById('b-p'),idleT;
function show(i){cur=(i+N)%N;
 phs.forEach(function(p,j){p.classList.toggle('active',j===cur)});
 cnt.textContent=(cur+1)+' / '+N;
 bar.classList.remove('run');void bar.offsetWidth;bar.classList.add('run');
 armTimer(DUR)}
function armTimer(ms){clearTimeout(timer);left=ms;t0=Date.now();
 if(playing)timer=setTimeout(step,ms)}
function step(){if(cur===N-1&&!LOOP){pause();return}show(cur+1)}
function pause(){playing=false;clearTimeout(timer);
 left=Math.max(0,left-(Date.now()-t0));document.body.classList.add('paused');
 btnP.textContent='▶'}
function play(){playing=true;document.body.classList.remove('paused');
 btnP.textContent='⏸';t0=Date.now();clearTimeout(timer);
 timer=setTimeout(step,left>150?left:DUR)}
function toggle(){playing?pause():play()}
function fullscreen(){document.fullscreenElement?document.exitFullscreen():
 document.documentElement.requestFullscreen()}
addEventListener('keydown',function(e){
 if(e.key===' '){e.preventDefault();toggle()}
 else if(e.key==='ArrowRight'||e.key==='PageDown')show(cur+1);
 else if(e.key==='ArrowLeft'||e.key==='PageUp')show(cur-1);
 else if(e.key==='Home')show(0);else if(e.key==='End')show(N-1);
 else if(e.key==='f'||e.key==='F')fullscreen()});
phs.forEach(function(p){p.addEventListener('click',toggle)});
document.getElementById('b-prev').onclick=function(e){e.stopPropagation();show(cur-1)};
document.getElementById('b-next').onclick=function(e){e.stopPropagation();show(cur+1)};
btnP.onclick=function(e){e.stopPropagation();toggle()};
document.getElementById('b-f').onclick=function(e){e.stopPropagation();fullscreen()};
var tx=0,ty=0;
addEventListener('touchstart',function(e){tx=e.touches[0].clientX;ty=e.touches[0].clientY},{passive:true});
addEventListener('touchend',function(e){var dx=e.changedTouches[0].clientX-tx,
 dy=e.changedTouches[0].clientY-ty;
 if(Math.abs(dx)>50&&Math.abs(dx)>Math.abs(dy))show(cur+(dx<0?1:-1))},{passive:true});
addEventListener('mousemove',function(){document.body.classList.remove('idle');
 clearTimeout(idleT);idleT=setTimeout(function(){document.body.classList.add('idle')},2600)});
show(0);
setTimeout(function(){document.body.classList.add('idle')},2600);
})();
"""


def export_slideshow_html(
    folder: str | Path,
    out: Optional[str | Path] = None,
    theme: str = "gallery-black",
    duration: float = 5.0,
    zoom: int = 12,
    title: str = "",
    subtitle: str = "",
    order: str = "name",
    captions: bool = False,
    intro: Optional[bool] = None,      # None = tu dong (>= 4 anh)
    loop: bool = True,
    size: str = "16:9",
    overwrite: bool = False,
    status: Optional[Callable[[str], None]] = None,
) -> tuple[Path, list[str]]:
    """Thu muc anh -> file HTML trinh chieu Ken Burns. Tra ve (path, canh bao)."""
    import html as _h

    status = status or (lambda s: None)
    t = get_slide_theme(theme)
    folder, files, warnings = _collect_photos(folder, order)
    if intro is None:
        intro = len(files) >= 4

    slides: list[str] = []
    if title:
        sub = f"<p>{_h.escape(subtitle)}</p>" if subtitle else ""
        slides.append(
            f'<div class="ph title"><div class="tt"><h1>{_h.escape(title)}</h1>'
            f'{sub}<p class="hint">Space: dừng/chạy · ←/→: chuyển ảnh · '
            f"F: toàn màn hình</p></div></div>")

    if intro:
        img = _intro_collage(folder, theme, size, status)
        if img is not None:
            uri = data_uri(pil_to_stream(img).getvalue(), "image/jpeg")
            slides.append(f'<div class="ph k1"><img src="{uri}" alt=""></div>')

    total = len(files)
    for i, f in enumerate(files):
        status(f"Đang nhúng ảnh {i + 1}/{total}...")
        try:
            uri = data_uri(*image_bytes_for_html(f, max_px=2000))
        except Exception as e:  # noqa: BLE001
            warnings.append(f"Bo qua {f.name}: {e}")
            continue
        kb = _KB_CLASSES[i % len(_KB_CLASSES)]
        cap = (f'<div class="cap">{_h.escape(f.stem)}</div>'
               if captions else "")
        slides.append(f'<div class="ph {kb}"><img src="{uri}" alt="">{cap}</div>')

    if not slides:
        raise CollageError("Khong nhung duoc anh nao.")

    html_doc = f"""<!DOCTYPE html>
<html lang="vi">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{_h.escape(title or folder.name)}</title>
<style>{_slideshow_css(t, duration, zoom, size)}</style>
</head>
<body>
{"".join(slides)}
<div id="bar"></div>
<div id="ui">
<span id="cnt"></span>
<button id="b-prev" title="Ảnh trước (←)">‹</button>
<button id="b-p" title="Dừng/chạy (Space)">⏸</button>
<button id="b-next" title="Ảnh sau (→)">›</button>
<button id="b-f" title="Toàn màn hình (F)">⛶</button>
</div>
<script>window.__DUR={max(1.5, duration)};window.__LOOP={str(bool(loop)).lower()};</script>
<script>{_SHOW_JS}</script>
</body>
</html>"""

    out_path = Path(out) if out else folder / "slideshow.html"
    if out_path.suffix.lower() not in (".html", ".htm"):
        out_path = out_path.with_suffix(".html")
    if not overwrite:
        out_path = unique_path(out_path)
    out_path.write_text(html_doc, encoding="utf-8")
    return out_path, warnings


# ================================================================== PPTX ====
_NS = ('xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
       'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
       'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"')

_TRANSITION_XML = '<p:transition {ns} spd="slow" advClick="1" advTm="{adv}"><p:fade/></p:transition>'

# Hieu ung Grow/Shrink (presetID 6) chay cung luc khi slide hien ra
_TIMING_XML = """<p:timing {ns}><p:tnLst><p:par>
<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
<p:seq concurrent="1" nextAc="seek">
<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
<p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="4" fill="hold">
<p:stCondLst><p:cond delay="0"/></p:stCondLst>
<p:childTnLst><p:par>
<p:cTn id="5" presetID="6" presetClass="emph" presetSubtype="0" dur="{dur}" fill="hold" nodeType="withEffect">
<p:stCondLst><p:cond delay="0"/></p:stCondLst>
<p:childTnLst><p:animScale><p:cBhvr>
<p:cTn id="6" dur="{dur}" fill="hold"/>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
</p:cBhvr><p:from x="{fx}" y="{fx}"/><p:to x="{tx}" y="{tx}"/></p:animScale>
</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
</p:cTn></p:par></p:childTnLst></p:cTn>
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""


def _inject_effects(slide, pic_shape_id: Optional[int], advance_ms: int,
                    dur_ms: int, zoom_in: bool, zoom: int) -> bool:
    """Chen XML chuyen canh fade + hieu ung zoom vao slide. True neu thanh cong."""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        sld = slide._element
        # thu tu schema: cSld, clrMapOvr, transition, timing
        anchor = sld.find(qn("p:clrMapOvr"))
        idx = list(sld).index(anchor) + 1 if anchor is not None else len(list(sld))

        trans = etree.fromstring(
            _TRANSITION_XML.format(ns=_NS, adv=advance_ms))
        sld.insert(idx, trans)

        if pic_shape_id is not None:
            k = 100000 + max(4, min(40, zoom)) * 1000
            fx, tx = (100000, k) if zoom_in else (k, 100000)
            timing = etree.fromstring(_TIMING_XML.format(
                ns=_NS, dur=dur_ms, spid=pic_shape_id, fx=fx, tx=tx))
            sld.insert(idx + 1, timing)
        return True
    except Exception:  # noqa: BLE001
        return False


def _cover_crop(img: Image.Image, tw: int, th: int) -> Image.Image:
    k = max(tw / img.width, th / img.height)
    nw, nh = round(img.width * k), round(img.height * k)
    img = img.resize((nw, nh), Image.Resampling.LANCZOS)
    x, y = (nw - tw) // 2, (nh - th) // 2
    return img.crop((x, y, x + tw, y + th))


def export_slideshow_pptx(
    folder: str | Path,
    out: Optional[str | Path] = None,
    theme: str = "gallery-black",
    duration: float = 5.0,
    zoom: int = 12,
    title: str = "",
    subtitle: str = "",
    order: str = "name",
    intro: Optional[bool] = None,
    size: str = "16:9",
    overwrite: bool = False,
    status: Optional[Callable[[str], None]] = None,
) -> tuple[Path, list[str]]:
    """Thu muc anh -> file .pptx: moi anh 1 slide full man hinh, zoom + fade
    + tu chuyen slide sau `duration` giay. Tra ve (path, canh bao)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise CollageError(
            "Chua cai thu vien python-pptx. Chay:  pip install python-pptx"
        ) from None

    status = status or (lambda s: None)
    t = get_slide_theme(theme)
    folder, files, warnings = _collect_photos(folder, order)
    if intro is None:
        intro = len(files) >= 4

    prs = Presentation()
    if size == "16:9":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    sw_in, sh_in = prs.slide_width.inches, prs.slide_height.inches
    tw, th = round(sw_in * 144), round(sh_in * 144)     # render 144 dpi
    blank = prs.slide_layouts[6]
    adv_ms = round(max(1.5, duration) * 1000)
    dur_ms = adv_ms
    xml_fail = False

    def _rgb(color: str):
        from pptx.dml.color import RGBColor
        return RGBColor.from_string(color.lstrip("#"))

    def _bg(slide):
        fill = slide.background.fill
        try:
            fill.gradient()
            stops = fill.gradient_stops
            stops[0].color.rgb = _rgb(t["bg"][0])
            stops[1].color.rgb = _rgb(t["bg"][1])
        except Exception:  # noqa: BLE001
            fill.solid()
            fill.fore_color.rgb = _rgb(t["bg"][0])

    # --- slide tieu de ---
    if title:
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        box = slide.shapes.add_textbox(
            Inches(sw_in * 0.1), Inches(sh_in * 0.32),
            Inches(sw_in * 0.8), Inches(sh_in * 0.36))
        tf = box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import PP_ALIGN
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = title
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = _rgb(t["h1"])
        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(20)
            r2.font.color.rgb = _rgb(t["muted"])
        if not _inject_effects(slide, None, adv_ms, dur_ms, True, zoom):
            xml_fail = True

    # --- slide collage mo dau ---
    photos: list[tuple[Image.Image, bool]] = []
    if intro:
        img = _intro_collage(folder, theme, size, status)
        if img is not None:
            photos.append((img, False))                 # zoom-out

    total = len(files)
    for i, f in enumerate(files):
        status(f"Đang dựng slide ảnh {i + 1}/{total} (PPTX)...")
        try:
            with Image.open(f) as im:
                im = ImageOps.exif_transpose(im).copy()
            photos.append((im, i % 2 == 0))
        except Exception as e:  # noqa: BLE001
            warnings.append(f"Bo qua {f.name}: {e}")

    if not photos:
        raise CollageError("Khong doc duoc anh nao.")

    for img, zoom_in in photos:
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        frame = _cover_crop(img, tw, th)
        pic = slide.shapes.add_picture(
            pil_to_stream(frame), 0, 0, prs.slide_width, prs.slide_height)
        if not _inject_effects(slide, pic.shape_id, adv_ms, dur_ms,
                               zoom_in, zoom):
            xml_fail = True

    if xml_fail:
        warnings.append(
            "Khong chen duoc hieu ung zoom/fade tu dong — file van dung, "
            "them hieu ung thu cong trong PowerPoint neu can.")
    else:
        warnings.append(
            "PPTX da co zoom + fade + tu chuyen slide. Muon lap vo han: "
            "PowerPoint > Slide Show > Set Up Slide Show > Loop continuously.")

    if title:
        prs.core_properties.title = title

    out_path = Path(out) if out else folder / "slideshow.pptx"
    if out_path.suffix.lower() != ".pptx":
        out_path = out_path.with_suffix(".pptx")
    if not overwrite:
        out_path = unique_path(out_path)
    prs.save(str(out_path))
    return out_path, warnings
